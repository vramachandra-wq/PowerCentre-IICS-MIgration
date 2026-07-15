"""Generate executive summary PowerPoint slide for PC to IICS migration."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent / "PowerCenter_IICS_Migration_Executive_Slide.pptx"

# Colors
BG_DARK = RGBColor(15, 23, 42)
BG_CARD = RGBColor(30, 41, 59)
ACCENT_CYAN = RGBColor(56, 189, 248)
ACCENT_PURPLE = RGBColor(167, 139, 250)
TEXT_WHITE = RGBColor(248, 250, 252)
TEXT_LIGHT = RGBColor(203, 213, 225)
TEXT_MUTED = RGBColor(148, 163, 184)
BLUE_CORE = RGBColor(37, 99, 235)
PURPLE_BUFFER = RGBColor(124, 58, 237)
BORDER = RGBColor(71, 85, 105)


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 12,
    bold: bool = False,
    color: RGBColor = TEXT_LIGHT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return box


def add_card(slide, left, top, width, height, title: str, title_color: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(card, BG_CARD)
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.75)

    add_textbox(
        slide,
        left + Inches(0.12),
        top + Inches(0.1),
        width - Inches(0.24),
        Inches(0.35),
        title,
        font_size=11,
        bold=True,
        color=title_color,
    )


def add_bullet_block(slide, left, top, width, height, bullets: list[str], font_size: int = 9) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {bullet}"
        p.space_after = Pt(4)
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = TEXT_LIGHT
        run.font.name = "Segoe UI"


def add_metric_box(slide, left, top, width, height, value: str, label: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(shape, RGBColor(30, 58, 90))
    shape.line.color.rgb = ACCENT_CYAN
    shape.line.width = Pt(0.5)

    add_textbox(
        slide, left, top + Inches(0.05), width, Inches(0.3),
        value, font_size=16, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, left, top + Inches(0.28), width, Inches(0.2),
        label.upper(), font_size=7, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )


def add_timeline_bar(slide, left, top, width, height, label: str, duration: str, bar_color: RGBColor, bar_text: str) -> None:
    add_textbox(slide, left, top, width * 0.7, Inches(0.18), label, font_size=8, color=TEXT_MUTED)
    add_textbox(slide, left + width * 0.7, top, width * 0.3, Inches(0.18), duration, font_size=8, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.2), width, height)
    set_fill(bar, bar_color)
    bar.line.fill.background()

    add_textbox(
        slide, left, top + Inches(0.2), width, height,
        bar_text, font_size=8, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
    )


def build_slide() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, BG_DARK)
    bg.line.fill.background()

    # Header
    add_textbox(
        slide, Inches(0.45), Inches(0.3), Inches(12.4), Inches(0.45),
        "PowerCenter → IICS / IDMC Migration Accelerator",
        font_size=24, bold=True, color=TEXT_WHITE,
    )
    add_textbox(
        slide, Inches(0.45), Inches(0.72), Inches(12.4), Inches(0.3),
        "Enterprise metadata migration platform — pilot proven, enterprise scale planned",
        font_size=12, color=TEXT_MUTED,
    )

    col_top = Inches(1.15)
    col_h = Inches(5.35)
    gap = Inches(0.2)
    col1_w = Inches(4.15)
    col2_w = Inches(4.15)
    col3_w = Inches(3.35)
    col1_l = Inches(0.45)
    col2_l = col1_l + col1_w + gap
    col3_l = col2_l + col2_w + gap

    # Column 1 — Proven Delivery
    add_card(slide, col1_l, col_top, col1_w, col_h, "1 · PROVEN DELIVERY (PILOT)", ACCENT_CYAN)
    add_bullet_block(
        slide, col1_l + Inches(0.15), col_top + Inches(0.48), col1_w - Inches(0.3), Inches(2.6),
        [
            "End-to-end pipeline: PC XML → canonical metadata → validation → remediation → IICS export",
            "Rule-based auto-remediation with remediated XML & import-ready IDMC packages",
            "Enterprise migration reports, readiness scores & consolidated dashboards",
            "AI-assisted validation & recommendations (FastAPI + Streamlit)",
            "MySQL central repository for metadata governance & future PC-vs-IICS comparison",
        ],
        font_size=9,
    )
    m_top = col_top + Inches(3.2)
    m_w = Inches(1.85)
    m_h = Inches(0.55)
    m_gap = Inches(0.15)
    metrics = [("98.5%", "Validation Pass"), ("99.3%", "Auto-Fix Rate"), ("92%", "Readiness Score"), ("73", "IICS Assets Exported")]
    for idx, (val, lbl) in enumerate(metrics):
        row, col = divmod(idx, 2)
        x = col1_l + Inches(0.15) + col * (m_w + m_gap)
        y = m_top + row * (m_h + Inches(0.1))
        add_metric_box(slide, x, y, m_w, m_h, val, lbl)

    # Column 2 — Long-Run Services
    add_card(slide, col2_l, col_top, col2_w, col_h, "2 · LONG-RUN SERVICES (6,000–7,000 ASSETS)", ACCENT_CYAN)
    add_bullet_block(
        slide, col2_l + Inches(0.15), col_top + Inches(0.48), col2_w - Inches(0.3), Inches(4.7),
        [
            "Batch metadata extraction — workflows, mappings, transformations at enterprise volume",
            "Scaled validation & remediation — datatype harmonization, rule engine, XML auto-fix",
            "IICS / IDMC package factory — bulk DTEMPLATE, MTT, TASKFLOW generation & import validation",
            "Migration intelligence — complexity scoring, risk assessment, executive dashboards",
            "AI remediation studio — complex patterns (Stored Proc, Java, CDC, SAP) with guided fixes",
            "Lineage & PC-vs-IICS comparison — central repository diff, drift tracking, cutover readiness",
            "CI/CD & automation ops — scheduled runs, GitHub Actions, API-driven migration factory",
        ],
        font_size=8.5,
    )

    # Column 3 — Timeline
    add_card(slide, col3_l, col_top, col3_w, col_h, "3 · PROGRAM TIMELINE", ACCENT_PURPLE)
    t_left = col3_l + Inches(0.15)
    t_width = col3_w - Inches(0.3)
    t_y = col_top + Inches(0.55)

    add_timeline_bar(
        slide, t_left, t_y, t_width, Inches(0.35),
        "Core delivery", "9 months", BLUE_CORE, "Build · Migrate · Validate · Export",
    )
    add_timeline_bar(
        slide, t_left, t_y + Inches(0.75), t_width, Inches(0.35),
        "Risk buffer", "2 months", PURPLE_BUFFER, "Complex assets · Regression · Cutover",
    )

    # Combined bar
    add_textbox(slide, t_left, t_y + Inches(1.45), t_width * 0.7, Inches(0.18), "Total program", font_size=8, color=TEXT_MUTED)
    add_textbox(slide, t_left + t_width * 0.7, t_y + Inches(1.45), t_width * 0.3, Inches(0.18), "11 months", font_size=8, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

    core_w = t_width * (9 / 11)
    buffer_w = t_width * (2 / 11)
    bar_y = t_y + Inches(1.65)
    core_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, t_left, bar_y, core_w, Inches(0.35))
    set_fill(core_bar, BLUE_CORE)
    core_bar.line.fill.background()
    add_textbox(slide, t_left, bar_y, core_w, Inches(0.35), "9 mo", font_size=9, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    buffer_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, t_left + core_w, bar_y, buffer_w, Inches(0.35))
    set_fill(buffer_bar, PURPLE_BUFFER)
    buffer_bar.line.fill.background()
    add_textbox(slide, t_left + core_w, bar_y, buffer_w, Inches(0.35), "+2 mo", font_size=8, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_textbox(
        slide, t_left, t_y + Inches(2.25), t_width, Inches(0.4),
        "11 Months", font_size=22, bold=True, color=ACCENT_PURPLE, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, t_left, t_y + Inches(2.65), t_width, Inches(0.25),
        "9 months execution + 2 months buffer", font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )

    # Footer
    footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(6.65), Inches(12.4), Pt(1))
    set_fill(footer_line, BORDER)
    footer_line.line.fill.background()

    add_textbox(
        slide, Inches(0.45), Inches(6.75), Inches(9.5), Inches(0.3),
        "Pilot scope: 14 XML files · 23 mappings · 202 canonical assets · Custom_Project IDMC export",
        font_size=9, color=TEXT_MUTED,
    )

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(6.72), Inches(2.65), Inches(0.35))
    set_fill(badge, RGBColor(22, 78, 99))
    badge.line.color.rgb = ACCENT_CYAN
    badge.line.width = Pt(0.5)
    add_textbox(
        slide, Inches(10.2), Inches(6.72), Inches(2.65), Inches(0.35),
        "Platform Ready for Enterprise Scale-Up",
        font_size=8, bold=True, color=RGBColor(103, 232, 249), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_slide()
    print(f"Created: {path}")
